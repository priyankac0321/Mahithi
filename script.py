from os import listdir
from os.path import isfile, join,isdir

import sys

from langchain_huggingface import HuggingFaceEmbeddings
from langchain_qdrant import Qdrant
from langchain_text_splitters import TokenTextSplitter

from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams

from PyPDF2 import PdfReader
from pptx import Presentation
import docx



    
embedding = HuggingFaceEmbeddings(model_name = "sentence-transformers/all-mpnet-base-v2",
                           model_kwargs = {'device':'cuda'},
                           encode_kwargs = {'normalize_embeddings':True} ) 
    
client = QdrantClient(path="Database/")
collection_name = "db"

if client.collection_exists(collection_name):
        client.delete_collection(collection_name)


client.create_collection(collection_name,vectors_config=VectorParams(size=768, distance=Distance.COSINE))
qdrant = Qdrant(client, collection_name, embedding)



def getfiles(dir):
    
    file_list = []
    for f in listdir(dir):
        if isfile(join(dir,f)):
            file_list.append(join(dir,f))
        elif isdir(join(dir,f)):
            file_list += getfiles(join(dir,f))

    return file_list



def wordfiles(filename):
    
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)


def pptfiles(filename):
    
    fullText = []
    prs = Presentation(filename)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"): 
                fullText.append(shape.text)
    return "\n".join(fullText)


def chunker(content):

    text_splitter = TokenTextSplitter(chunk_size=500, chunk_overlap=100)
    chunks = text_splitter.split_text(content)
    
    return chunks


def extract(mypath):

    onlyfiles = getfiles(mypath)
    file_content = ""

    for file in onlyfiles:
        file_content = ""

        if file.endswith(".pdf"):
            print("indexing "+file)
            reader = PdfReader(file)
            for i in range(0,len(reader.pages)):
                file_content = file_content + " "+reader.pages[i].extract_text()

        elif file.endswith(".docx"):
            print("indexing " + file)
            file_content = wordfiles(file)

        elif file.endswith(".pptx"):
            print("indexing " + file)
            file_content = pptfiles(file)

        elif file.endswith(".txt") or file.endswith(".md") or file.endswith(".markdown"):
            print("indexing " + file)
            f = open(file,'r')
            file_content = f.read()
            f.close()

        else:
            continue

        chunk = chunker(file_content)
        
        metadata = []

        for i in range(0,len(chunk)):
            metadata.append({"path":file})
        qdrant.add_texts(chunk , metadatas=metadata)
        len(chunk)

    print(onlyfiles)
       



if __name__ == "__main__":
    arguments = sys.argv
    if len(arguments)>1:
        extract(arguments[1])
    else:
        print("You need to provide a path to folder with documents to index as command line argument")